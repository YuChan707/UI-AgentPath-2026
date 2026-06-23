"""
embedding-service
  Subscribes: document-submitted
  Publishes:  embedding-complete, pipeline-status

Pipeline step 1 of 4:
  download file bytes from PostgreSQL
  → extract text (PDF/PPTX/DOCX/TXT)
  → clean text
  → chunk by page / paragraph
  → embed with Ollama (nomic-embed-text)
  → store vectors in ChromaDB
  → publish embedding-complete
"""
import io
import json
import logging
import os
import re
import uuid
import asyncio
from contextlib import asynccontextmanager
from typing import Any

import asyncpg
import chromadb
import httpx
import ollama
import pypdf
from docx import Document as DocxDocument
from fastapi import FastAPI, Request
from fastapi.responses import JSONResponse
from pptx import Presentation

logging.basicConfig(level=logging.INFO)
log = logging.getLogger("embedding-service")

DATABASE_URL = os.getenv("DATABASE_URL", "postgresql://postgres:postgres@localhost:5432/onlooker")
CHROMA_HOST  = os.getenv("CHROMA_HOST", "localhost")
CHROMA_PORT  = int(os.getenv("CHROMA_PORT", "8000"))
OLLAMA_HOST  = os.getenv("OLLAMA_HOST", "http://localhost:11434")
EMBED_MODEL  = os.getenv("EMBED_MODEL", "nomic-embed-text")
DAPR_HTTP    = f"http://localhost:{os.getenv('DAPR_HTTP_PORT', '3500')}"
PUBSUB_NAME  = "onlooker-pubsub"

_pool: asyncpg.Pool | None = None
_chroma: chromadb.HttpClient | None = None


@asynccontextmanager
async def lifespan(app: FastAPI):
    global _pool, _chroma
    # Strip asyncpg dialect prefix that SQLAlchemy needs but asyncpg rejects
    dsn = DATABASE_URL.replace("postgresql+asyncpg://", "postgresql://")
    _pool   = await asyncpg.create_pool(dsn, min_size=2, max_size=10)
    _chroma = chromadb.HttpClient(host=CHROMA_HOST, port=CHROMA_PORT)
    log.info("embedding-service ready")
    yield
    await _pool.close()


app = FastAPI(title="embedding-service", lifespan=lifespan)


# ── Dapr subscription registration ───────────────────────────────────────────

@app.get("/dapr/subscribe")
def dapr_subscribe():
    return [{"pubsubname": PUBSUB_NAME, "topic": "document-submitted", "route": "/document-submitted"}]


# ── Internal helpers ─────────────────────────────────────────────────────────

async def _publish(topic: str, data: dict[str, Any]) -> None:
    url = f"{DAPR_HTTP}/v1.0/publish/{PUBSUB_NAME}/{topic}"
    async with httpx.AsyncClient() as client:
        await client.post(url, json=data, timeout=30)


async def _status(run_id: str, stage: str, detail: str = "") -> None:
    await _publish("pipeline-status", {"run_id": run_id, "stage": stage, "detail": detail})


def _clean(text: str) -> str:
    text = re.sub(r"/\*.*?\*/", " ", text, flags=re.DOTALL)
    text = re.sub(r'[()#"\']+', " ", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _extract_pages(file_bytes: bytes, ext: str) -> list[str]:
    """Return one string per logical page/slide/chunk."""
    ext = ext.lower().lstrip(".")
    pages: list[str] = []

    if ext == "pdf":
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        for page in reader.pages:
            pages.append(page.extract_text() or "")

    elif ext in ("pptx", "ppt"):
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            text = " ".join(
                shape.text_frame.text
                for shape in slide.shapes
                if shape.has_text_frame
            )
            pages.append(text)

    elif ext in ("docx", "doc"):
        doc = DocxDocument(io.BytesIO(file_bytes))
        full = "\n".join(p.text for p in doc.paragraphs)
        words = full.split()
        for i in range(0, len(words), 500):
            pages.append(" ".join(words[i : i + 500]))

    else:  # txt, md, csv, …
        raw = file_bytes.decode("utf-8", errors="ignore")
        words = raw.split()
        for i in range(0, len(words), 500):
            pages.append(" ".join(words[i : i + 500]))

    return [_clean(p) for p in pages if _clean(p)]


def _embed_sync(chunks: list[str]) -> list[list[float]]:
    client = ollama.Client(host=OLLAMA_HOST)
    return [
        client.embeddings(model=EMBED_MODEL, prompt=chunk)["embedding"]
        for chunk in chunks
    ]


# ── Event handler ─────────────────────────────────────────────────────────────

@app.post("/document-submitted")
async def on_document_submitted(request: Request) -> JSONResponse:
    body        = await request.json()
    data        = body.get("data", body)
    run_id      = data.get("run_id") or str(uuid.uuid4())
    document_id = data["document_id"]
    file_ext    = data.get("file_extension", "pdf")

    log.info("run=%s doc=%s", run_id, document_id)
    await _status(run_id, "embedding:started", "Downloading file from database")

    # 1 — Fetch raw bytes from PostgreSQL
    async with _pool.acquire() as conn:
        row = await conn.fetchrow(
            "SELECT file_data FROM documents WHERE document_id = $1", uuid.UUID(document_id)
        )
    if not row:
        log.error("Document %s not found", document_id)
        return JSONResponse({"status": "SUCCESS"})

    file_bytes = bytes(row["file_data"])

    # 2 — Extract + clean text
    await _status(run_id, "embedding:extracting", "Extracting text")
    chunks = _extract_pages(file_bytes, file_ext)
    if not chunks:
        log.warning("No text extracted from %s", document_id)
        return JSONResponse({"status": "SUCCESS"})

    # 3 — Embed with Ollama (blocking, run in thread pool)
    await _status(run_id, "embedding:vectorizing", f"Vectorizing {len(chunks)} chunks")
    loop   = asyncio.get_event_loop()
    embeds = await loop.run_in_executor(None, _embed_sync, chunks)

    # 4 — Store in ChromaDB
    col_name   = f"run_{run_id.replace('-', '_')}"
    collection = _chroma.create_collection(col_name, get_or_create=True)
    collection.add(
        documents=chunks,
        embeddings=embeds,
        ids=[f"page_{i}" for i in range(len(chunks))],
        metadatas=[{"page": i, "run_id": run_id} for i in range(len(chunks))],
    )

    # 5 — Persist run record
    async with _pool.acquire() as conn:
        await conn.execute(
            """
            INSERT INTO pipeline_runs
              (run_id, document_id, collection_name, status, audience_config, insight_flags)
            VALUES ($1, $2, $3, 'embedding_complete', $4, $5)
            ON CONFLICT (run_id) DO UPDATE
              SET collection_name = $3, status = 'embedding_complete'
            """,
            uuid.UUID(run_id),
            uuid.UUID(document_id),
            col_name,
            json.dumps(data.get("audience_config", {})),
            json.dumps(data.get("insight_flags", {})),
        )

    # 6 — Publish to next stage
    await _status(run_id, "embedding:complete", f"Stored {len(chunks)} vectors")
    await _publish("embedding-complete", {
        "run_id":         run_id,
        "collection_name": col_name,
        "chunk_count":    len(chunks),
        "audience_config": data.get("audience_config", {}),
        "insight_flags":  data.get("insight_flags", {}),
    })

    return JSONResponse({"status": "SUCCESS"})
