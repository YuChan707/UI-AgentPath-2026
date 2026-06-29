import io
import docx
from fastapi import UploadFile, HTTPException
from pptx import Presentation
from pypdf import PdfReader


async def extract_text(file: UploadFile) -> str:
    """Extract plain text from .pptx, .docx, or .pdf uploads."""
    content = await file.read()
    name = (file.filename or "").lower()

    if name.endswith(".pptx"):
        return _from_pptx(content)
    if name.endswith(".docx"):
        return _from_docx(content)
    if name.endswith(".pdf"):
        return _from_pdf(content)

    raise HTTPException(
        status_code=415,
        detail=f"Unsupported file type. Accepted: .pptx, .docx, .pdf"
    )


def _from_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
    return "\n".join(parts)


def _from_docx(content: bytes) -> str:
    doc = docx.Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _from_pdf(content: bytes) -> str:
    reader = PdfReader(io.BytesIO(content))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages)
