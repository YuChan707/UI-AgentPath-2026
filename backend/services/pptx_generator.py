import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BLUE   = RGBColor(0x00, 0x78, 0xD4)
GREEN  = RGBColor(0x10, 0x79, 0x3F)
PURPLE = RGBColor(0x8B, 0x00, 0x8B)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT  = RGBColor(0xF5, 0xF5, 0xF5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _header_bar(slide, prs, color, title):
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12), Inches(0.5))
    tf = tb.text_frame
    tf.text = title
    tf.paragraphs[0].runs[0].font.color.rgb = WHITE
    tf.paragraphs[0].runs[0].font.size = Pt(20)
    tf.paragraphs[0].runs[0].font.bold = True

def _text_box(slide, left, top, width, height, text, size=12, bold=False, color=DARK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def generate_onlooker_report(session_data: dict) -> bytes:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    a = session_data.get("analytics", {})

    # Slide 1: Title
    slide = _blank_slide(prs)
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    _text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
              "Onlooker — Audience Intelligence Report", 30, True, WHITE)
    _text_box(slide, Inches(0.5), Inches(0.9), Inches(12), Inches(0.4),
              f"Persona: {session_data.get('persona_type','').title()}  |  "
              f"Region: {session_data.get('region','').upper()}  |  "
              f"Focus: {session_data.get('focus_area','').title()}", 13, False, WHITE)

    # Slide 2: Analytics
    slide = _blank_slide(prs)
    _header_bar(slide, prs, BLUE, "Audience Analytics")
    metrics = [
        ("Engagement",   a.get("engagement_level", 0), BLUE),
        ("Conviction",   a.get("conviction_level", 0),  GREEN),
        ("Clarity",      a.get("clarity_score", 0),     PURPLE),
    ]
    for i, (label, score, color) in enumerate(metrics):
        top = Inches(1.0 + i * 1.8)
        _text_box(slide, Inches(0.5), top, Inches(4), Inches(0.4), label, 14, True)
        pct = min(100, int(score * 100))
        bg = slide.shapes.add_shape(1, Inches(0.5), top + Inches(0.5), Inches(10), Inches(0.35))
        bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT; bg.line.color.rgb = color
        fill_w = max(Inches(0.1), Inches(10 * score))
        fill = slide.shapes.add_shape(1, Inches(0.5), top + Inches(0.5), fill_w, Inches(0.35))
        fill.fill.solid(); fill.fill.fore_color.rgb = color; fill.line.fill.background()
        _text_box(slide, Inches(10.8), top + Inches(0.45), Inches(1.5), Inches(0.35), f"{pct}%", 13, True)

    # Slide 3: Speech Metrics
    slide = _blank_slide(prs)
    _header_bar(slide, prs, GREEN, "Speech Performance")
    stats = [
        ("Pace (WPM)",    str(int(a.get("pace_wpm", 0)))),
        ("Filler Words",  str(int(a.get("filler_count", 0)))),
        ("Clarity Score", f"{a.get('clarity_score', 0):.0%}"),
    ]
    for i, (label, value) in enumerate(stats):
        top = Inches(1.2 + i * 1.8)
        _text_box(slide, Inches(1), top, Inches(5), Inches(0.4), label, 14, True)
        _text_box(slide, Inches(6), top, Inches(4), Inches(0.4), value, 28, True, GREEN)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
