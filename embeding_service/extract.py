"""Extract, clean and page-chunk the raw document bytes.

One chunk per "page": a slide (pptx), a page (pdf), a form-feed/size block
(txt, md) or a size block (docx). The text is cleaned of stray symbols that
extraction tends to leave behind (/ * ( ) # " ' …) so the embeddings see prose,
not markup.
"""

from __future__ import annotations

import io
import re

# Symbols that extraction leaves as noise; collapse them to spaces.
_NOISE = re.compile(r"[/\*\(\)#\"'`_~|\\\[\]{}<>]+")
_WS = re.compile(r"[ \t\f\v]+")
_BLANK_LINES = re.compile(r"\n{3,}")

# Roughly one "page" worth of characters when the format has no real pages.
_PAGE_CHARS = 1800


def clean(text: str) -> str:
    """Strip stray extraction symbols and normalise whitespace."""
    text = _NOISE.sub(" ", text)
    text = _WS.sub(" ", text)
    text = _BLANK_LINES.sub("\n\n", text)
    return text.strip()


def _by_size(text: str, size: int = _PAGE_CHARS) -> list[str]:
    """Split a long text into ~size-char blocks on paragraph boundaries."""
    blocks: list[str] = []
    buff: list[str] = []
    length = 0
    for para in re.split(r"\n\s*\n", text):
        para = para.strip()
        if not para:
            continue
        if length + len(para) > size and buff:
            blocks.append("\n\n".join(buff))
            buff, length = [], 0
        buff.append(para)
        length += len(para)
    if buff:
        blocks.append("\n\n".join(buff))
    return blocks


def _from_txt(content: bytes) -> list[str]:
    text = content.decode("utf-8", errors="replace")
    # Honour explicit page breaks (form feed) if the file has them.
    if "\f" in text:
        return [p for p in (s.strip() for s in text.split("\f")) if p]
    return _by_size(text)


def _from_pdf(content: bytes) -> list[str]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    return [(page.extract_text() or "").strip() for page in reader.pages]


def _from_pptx(content: bytes) -> list[str]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    pages: list[str] = []
    for slide in prs.slides:
        parts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        pages.append("\n".join(parts))
    return pages


def _from_docx(content: bytes) -> list[str]:
    import docx

    doc = docx.Document(io.BytesIO(content))
    text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return _by_size(text)


_EXTRACTORS = {
    "txt": _from_txt,
    "md": _from_txt,
    "pdf": _from_pdf,
    "pptx": _from_pptx,
    "docx": _from_docx,
    "doc": _from_docx,  # best effort; legacy .doc may fail to parse
}


def pages(content: bytes, doc_type: str) -> list[str]:
    """Return cleaned, non-empty page chunks for the given document type."""
    doc_type = (doc_type or "").lower().lstrip(".")
    extractor = _EXTRACTORS.get(doc_type)
    if extractor is None:
        raise ValueError(f"Unsupported doc_type: {doc_type!r}")
    raw_pages = extractor(content)
    cleaned = [clean(p) for p in raw_pages]
    return [p for p in cleaned if p]


__all__ = ["pages", "clean"]
